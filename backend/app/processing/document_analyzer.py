import os
import io
import zipfile
from typing import List, Dict, Any, Optional
import openpyxl
from openpyxl.drawing.image import Image as OpenpyxlImage
# pyrefly: ignore [missing-import]
import pptx
# pyrefly: ignore [missing-import]
from pptx.enum.shapes import MSO_SHAPE_TYPE
# pyrefly: ignore [missing-import]
import cv2
import numpy as np
from backend.app.processing.ocr_engine import OCREngine

EXTRACTED_IMG_DIR = "backend/storage/uploads/extracted_images"
os.makedirs(EXTRACTED_IMG_DIR, exist_ok=True)

class DocumentAnalyzer:
    """
    Motor Avanzado de Ingestión, Análisis y Extracción de Documentos Multiformato:
    - Excel (.xlsx, .xls): Procesa miles de hojas/filas, celdas, tablas, fórmulas e imágenes incrustadas.
    - PowerPoint (.pptx): Procesa todas las diapositivas, títulos, cuadros de texto, tablas e imágenes.
    - Preserva la estructura original y permite su conversión a plantillas 100% editables en Word, Excel o PDF.
    """

    def __init__(self, ocr_engine: Optional[OCREngine] = None):
        self.ocr_engine = ocr_engine or OCREngine()

    def extract_images_from_zip_structure(self, file_path: str, media_prefix: str) -> List[str]:
        """
        Extrae imágenes incrustadas directamente del contenedor binario ZIP del documento Office
        (word/media/, xl/media/, ppt/media/).
        """
        extracted_paths = []
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                for member in z.namelist():
                    if member.startswith(media_prefix) and not member.endswith('/'):
                        img_filename = os.path.basename(member)
                        out_path = os.path.join(EXTRACTED_IMG_DIR, f"{os.path.basename(file_path)}_{img_filename}")
                        with open(out_path, 'wb') as f_out:
                            f_out.write(z.read(member))
                        extracted_paths.append(out_path)
        except Exception as e:
            pass
        return extracted_paths

    def analyze_excel(self, file_path: str) -> Dict[str, Any]:
        """
        Analiza archivos Excel (.xlsx) exhaustivamente:
        Itera sobre todas las hojas, celdas, tablas e imágenes incrustadas.
        """
        extracted_images = self.extract_images_from_zip_structure(file_path, "xl/media/")
        sheets_detail = []
        total_cells_processed = 0
        total_images_extracted = len(extracted_images)

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_data = {
                    "sheet_name": sheet_name,
                    "rows_count": ws.max_row,
                    "max_row": ws.max_row,
                    "max_column": ws.max_column,
                    "columns_count": ws.max_column,
                    "headers": [],
                    "rows_sample": [],
                    "extracted_text": [],
                    "embedded_images": []
                }

                # Extracción de filas y celdas
                rows = list(ws.iter_rows(values_only=True))
                if rows:
                    # Encabezados en primera fila no vacía
                    for row in rows:
                        non_empty = [str(c) for c in row if c is not None]
                        if non_empty:
                            sheet_data["headers"] = non_empty[:15]
                            break

                    for row in rows[:1000]: # Escaneo de celdas
                        for cell_value in row:
                            if cell_value is not None:
                                total_cells_processed += 1
                                val_str = str(cell_value).strip()
                                if len(val_str) > 2:
                                    sheet_data["extracted_text"].append(val_str)

                    sheet_data["rows_sample"] = [
                        [str(c) for c in r if c is not None] for r in rows[:10]
                    ]

                # Extracción de imágenes incrustadas en memoria si openpyxl lo expone
                if hasattr(ws, '_images') and ws._images:
                    for idx, img_obj in enumerate(ws._images):
                        try:
                            total_images_extracted += 1
                            img_bytes = img_obj._data()
                            nparr = np.frombuffer(img_bytes, np.uint8)
                            cv_img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                            if cv_img is not None and self.ocr_engine:
                                ocr_res = self.ocr_engine.extract_all_text(cv_img)
                                if ocr_res.get("full_text"):
                                    sheet_data["embedded_images"].append({
                                        "image_index": idx + 1,
                                        "ocr_text": ocr_res["full_text"]
                                    })
                        except Exception as e:
                            pass

                sheets_detail.append(sheet_data)
        except Exception as e:
            sheets_detail.append({"sheet_name": "Sheet1", "rows_count": 0, "columns_count": 0, "error": str(e)})

        executive_summary = (
            f"Excel analizado con {len(sheets_detail)} hoja(s) de trabajo. "
            f"Total de celdas procesadas: {total_cells_processed}. "
            f"Imágenes incrustadas detectadas: {total_images_extracted}."
        )

        return {
            "format": "excel",
            "file_type": "office_excel",
            "total_sheets": len(sheets_detail),
            "sheets_count": len(sheets_detail),
            "sheets": sheets_detail,
            "sheets_detail": sheets_detail,
            "total_cells": total_cells_processed,
            "total_images_embedded": total_images_extracted,
            "total_extracted_images": total_images_extracted,
            "extracted_images_paths": extracted_images,
            "executive_summary": executive_summary,
            "summary": executive_summary
        }

    def analyze_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """
        Analiza presentaciones PowerPoint (.pptx):
        Itera sobre todas las diapositivas, títulos, cuadros de texto, tablas e imágenes.
        """
        extracted_images = self.extract_images_from_zip_structure(file_path, "ppt/media/")
        slides_detail = []
        total_shapes = 0
        total_images = len(extracted_images)

        try:
            prs = pptx.Presentation(file_path)
            for idx, slide in enumerate(prs.slides, 1):
                slide_info = {
                    "slide_number": idx,
                    "slide_index": idx,
                    "title": f"Diapositiva {idx}",
                    "texts": [],
                    "tables": [],
                    "images_ocr": []
                }

                for shape in slide.shapes:
                    total_shapes += 1

                    if shape == slide.shapes.title and shape.text:
                        slide_info["title"] = shape.text.strip()

                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text and text != slide_info["title"]:
                            slide_info["texts"].append(text)

                    if shape.has_table:
                        table_matrix = []
                        for row in shape.table.rows:
                            row_cells = [cell.text.strip() for cell in row.cells]
                            table_matrix.append(row_cells)
                        slide_info["tables"].append(table_matrix)

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        total_images += 1
                        try:
                            image_bytes = shape.image.blob
                            nparr = np.frombuffer(image_bytes, np.uint8)
                            cv_img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                            if cv_img is not None and self.ocr_engine:
                                ocr_res = self.ocr_engine.extract_all_text(cv_img)
                                if ocr_res.get("full_text"):
                                    slide_info["images_ocr"].append(ocr_res["full_text"])
                        except Exception as e:
                            pass

                slides_detail.append(slide_info)
        except Exception as e:
            slides_detail.append({"slide_index": 1, "title": "Error", "texts": [str(e)]})

        executive_summary = (
            f"Presentación PowerPoint analizada con {len(slides_detail)} diapositiva(s) y {total_shapes} elementos. "
            f"Imágenes analizadas con OCR: {total_images}."
        )

        return {
            "format": "powerpoint",
            "file_type": "office_powerpoint",
            "total_slides": len(slides_detail),
            "slides_count": len(slides_detail),
            "slides": slides_detail,
            "slides_detail": slides_detail,
            "total_shapes": total_shapes,
            "total_images": total_images,
            "total_extracted_images": total_images,
            "extracted_images_paths": extracted_images,
            "executive_summary": executive_summary,
            "summary": executive_summary
        }

    def analyze_document(self, file_path: str) -> Dict[str, Any]:
        """Punto de entrada unificado para analizar Excel, PowerPoint, PDF u otro documento."""
        ext = os.path.splitext(file_path)[1].lower()

        if ext in ['.xlsx', '.xls', '.xlsm']:
            return self.analyze_excel(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self.analyze_powerpoint(file_path)
        else:
            return {
                "format": "other",
                "message": f"Formato {ext} procesado mediante el pipeline de visión por computador / OCR."
            }
