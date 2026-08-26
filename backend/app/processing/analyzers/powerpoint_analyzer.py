import os
import zipfile
from typing import Dict, Any, List
import pptx

EXTRACTED_IMG_DIR = "backend/storage/uploads/extracted_images"
os.makedirs(EXTRACTED_IMG_DIR, exist_ok=True)

class PowerpointAnalyzer:
    """
    Analizador Especializado de Presentaciones Microsoft PowerPoint (.pptx).
    """

    def extract_embedded_images(self, file_path: str) -> List[str]:
        extracted_paths = []
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                for member in z.namelist():
                    if member.startswith('ppt/media/') and not member.endswith('/'):
                        img_filename = os.path.basename(member)
                        out_path = os.path.join(EXTRACTED_IMG_DIR, f"{os.path.basename(file_path)}_{img_filename}")
                        with open(out_path, 'wb') as f_out:
                            f_out.write(z.read(member))
                        extracted_paths.append(out_path)
        except Exception:
            pass
        return extracted_paths

    def analyze(self, file_path: str) -> Dict[str, Any]:
        size_bytes = os.path.getsize(file_path)
        extracted_images = self.extract_embedded_images(file_path)
        slides_detail = []

        try:
            prs = pptx.Presentation(file_path)
            for idx, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        texts.append(shape.text_frame.text.strip())

                slides_detail.append({
                    "slide_index": idx,
                    "slide_number": idx,
                    "title": texts[0] if texts else f"Diapositiva {idx}",
                    "texts": texts
                })
        except Exception as e:
            slides_detail.append({"slide_index": 1, "title": "Error", "texts": [str(e)]})

        executive_summary = (
            f"PowerPoint analizado con {len(slides_detail)} diapositiva(s). "
            f"Imágenes extraídas: {len(extracted_images)}."
        )

        return {
            "format": "powerpoint",
            "file_type": "office_powerpoint",
            "file_name": os.path.basename(file_path),
            "size_bytes": size_bytes,
            "total_slides": len(slides_detail),
            "slides_count": len(slides_detail),
            "slides": slides_detail,
            "slides_detail": slides_detail,
            "total_extracted_images": len(extracted_images),
            "extracted_images_paths": extracted_images,
            "executive_summary": executive_summary,
            "summary": executive_summary
        }
